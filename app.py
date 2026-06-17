import io
from datetime import datetime
from zoneinfo import ZoneInfo
from typing import Optional, List, Dict, Any

import numpy as np
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ============================================================
# PAGE CONFIG
# ============================================================
st.set_page_config(page_title="AI Agent - Delivery Manager", layout="wide")

# ============================================================
# BRAND COLORS
# ============================================================
PRIMARY = "#191D63"
DARK_BG = "#06172B"
CARD_BG = "#0D2742"
CARD_BG_2 = "#15263A"
TEXT = "#FFFFFF"
MUTED = "#D7E3F4"
CYAN = "#00D4FF"
GREEN = "#22C55E"
AMBER = "#F59E0B"
RED = "#EF4444"
PURPLE = "#8B5CF6"
BLUE = "#2F80ED"
GRID = "rgba(255,255,255,0.16)"

STATUS_COLOR_MAP = {
    "Done": PRIMARY,
    "Closed": PRIMARY,
    "Resolved": PRIMARY,
    "Completed": PRIMARY,
    "In Progress": GREEN,
    "Open": AMBER,
    "Blocked": RED,
    "Unknown": PURPLE,
}
PRIORITY_COLOR_MAP = {
    "Critical": RED,
    "High": AMBER,
    "Medium": BLUE,
    "Low": GREEN,
    "Unknown": PURPLE,
}
RISK_COLOR_MAP = {"Aging": AMBER, "Blocked": RED, "Normal": PRIMARY}

CHART_HEIGHT = 390

# ============================================================
# CSS
# ============================================================
def apply_app_style():
    st.markdown(
        """
        <style>
        :root {
          --primary: #191D63;
          --dark-bg: #06172B;
          --card-bg: #0D2742;
          --text: #FFFFFF;
          --muted: #D7E3F4;
          --cyan: #00D4FF;
          --green: #22C55E;
          --amber: #F59E0B;
          --red: #EF4444;
        }
        .stApp {
            background: linear-gradient(180deg, #F3F6FB 0%, #EEF3F9 100%);
        }
        .main .block-container {
            padding-top: 1.2rem;
            padding-left: 1.5rem;
            padding-right: 1.5rem;
            max-width: 100%;
        }
        .hero {
            background: linear-gradient(135deg, #06172B 0%, #0D2742 58%, #191D63 100%);
            color: white;
            border-radius: 22px;
            padding: 26px 30px;
            box-shadow: 0 18px 42px rgba(6, 23, 43, 0.22);
            margin-bottom: 18px;
        }
        .hero-top {
            display: flex;
            align-items: center;
            justify-content: space-between;
            gap: 20px;
        }
        .hero-badge {
            color: #00D4FF;
            font-weight: 800;
            letter-spacing: 0.08em;
            font-size: 0.86rem;
            text-transform: uppercase;
        }
        .hero h1 {
            font-size: 2.3rem;
            line-height: 1.1;
            margin: 10px 0 8px 0;
            color: white;
        }
        .hero p {
            color: #D7E3F4;
            font-size: 1.02rem;
            margin: 0;
        }
        .hero-meta {
            text-align: right;
            color: #D7E3F4;
            font-size: 0.92rem;
            min-width: 320px;
        }
        .upload-panel {
            background: #FFFFFF;
            border-radius: 18px;
            padding: 18px 22px;
            margin-bottom: 18px;
            border: 1px solid rgba(25, 29, 99, 0.12);
            box-shadow: 0 10px 26px rgba(15, 23, 42, 0.08);
        }
        .section-title {
            display: flex;
            align-items: center;
            gap: 0.55rem;
            font-size: 1.15rem;
            font-weight: 800;
            margin: 18px 0 10px 0;
            color: #0B1220;
        }
        .filter-note {
            background: #EAF2FF;
            color: #191D63;
            border-left: 5px solid #191D63;
            padding: 10px 14px;
            border-radius: 12px;
            font-weight: 600;
            margin: 8px 0 16px 0;
        }
        .kpi-grid {
            display: grid;
            grid-template-columns: repeat(5, minmax(170px, 1fr));
            gap: 14px;
            margin: 10px 0 16px 0;
        }
        .kpi-card {
            background: #0D2742;
            border: 1px solid rgba(255,255,255,0.14);
            border-radius: 18px;
            color: white;
            min-height: 124px;
            padding: 16px 18px;
            box-shadow: 0 12px 28px rgba(6, 23, 43, 0.20);
            position: relative;
            overflow: hidden;
        }
        .kpi-card::before {
            content: "";
            position: absolute;
            left: 0;
            top: 0;
            bottom: 0;
            width: 7px;
            background: var(--accent);
        }
        .kpi-title {
            font-size: 0.78rem;
            color: #D7E3F4;
            font-weight: 800;
            text-transform: uppercase;
            letter-spacing: 0.02em;
            margin-bottom: 18px;
        }
        .kpi-value {
            font-size: 2.0rem;
            font-weight: 900;
            line-height: 1.0;
            color: white;
            margin-bottom: 12px;
        }
        .kpi-subtitle {
            font-size: 0.78rem;
            color: #D7E3F4;
            line-height: 1.2;
        }
        .quality-panel {
            background: #06172B;
            color: #FFFFFF;
            border-radius: 18px;
            min-height: 440px;
            padding: 24px 26px;
            box-shadow: 0 12px 28px rgba(6, 23, 43, 0.18);
            border: 1px solid rgba(255,255,255,0.12);
            box-sizing: border-box;
            overflow: hidden;
        }
        .quality-eyebrow {
            color: #00D4FF;
            font-weight: 900;
            letter-spacing: 0.08em;
            text-transform: uppercase;
            font-size: 0.80rem;
            margin-bottom: 14px;
        }
        .quality-title {
            font-size: 1.15rem;
            font-weight: 900;
            margin-bottom: 10px;
        }
        .quality-value {
            font-size: 2.65rem;
            font-weight: 900;
            line-height: 1;
            margin-bottom: 10px;
        }
        .quality-status {
            display: inline-block;
            padding: 8px 14px;
            border-radius: 999px;
            font-weight: 900;
            font-size: 0.78rem;
            margin-bottom: 16px;
        }
        .quality-desc {
            color: #D7E3F4;
            font-size: 0.82rem;
            line-height: 1.38;
            margin-bottom: 14px;
        }
        .quality-mini-grid {
            display: grid;
            grid-template-columns: repeat(3, minmax(0, 1fr));
            gap: 9px;
        }
        .quality-mini-card {
            background: #15263A;
            border: 1px solid rgba(255,255,255,0.13);
            border-radius: 12px;
            padding: 10px 11px;
            min-height: 66px;
            box-sizing: border-box;
            overflow-wrap: anywhere;
        }
        .quality-mini-value {
            font-size: 1.02rem;
            font-weight: 900;
            margin-bottom: 5px;
            line-height: 1.2;
        }
        .quality-mini-label {
            color: #D7E3F4;
            font-size: 0.68rem;
            line-height: 1.2;
        }
        .hint {
            color: #64748B;
            font-size: 0.78rem;
            margin-top: 6px;
            margin-bottom: 14px;
        }

        html {
            scroll-behavior: smooth;
        }
        .app-top-anchor {
            display: block;
            position: relative;
            top: -16px;
            visibility: hidden;
        }
        .smooth-top-button {
            position: fixed;
            right: 28px;
            bottom: 28px;
            z-index: 999999;
            width: 54px;
            height: 54px;
            border-radius: 50%;
            background: linear-gradient(135deg, #191D63 0%, #00D4FF 100%);
            color: #FFFFFF !important;
            display: flex;
            align-items: center;
            justify-content: center;
            text-decoration: none !important;
            font-size: 1.35rem;
            font-weight: 900;
            box-shadow: 0 16px 34px rgba(25, 29, 99, 0.35);
            border: 1px solid rgba(255,255,255,0.35);
            transition: transform 0.22s ease, box-shadow 0.22s ease, opacity 0.22s ease;
        }
        .smooth-top-button:hover {
            transform: translateY(-5px) scale(1.04);
            box-shadow: 0 22px 44px rgba(25, 29, 99, 0.48);
            color: #FFFFFF !important;
            text-decoration: none !important;
        }
        .smooth-top-button::after {
            content: "Back To Top";
            position: absolute;
            right: 62px;
            top: 10px;
            background: #06172B;
            color: #FFFFFF;
            font-size: 0.72rem;
            font-weight: 800;
            letter-spacing: 0.02em;
            white-space: nowrap;
            padding: 8px 10px;
            border-radius: 10px;
            opacity: 0;
            pointer-events: none;
            transform: translateX(8px);
            transition: opacity 0.22s ease, transform 0.22s ease;
            border: 1px solid rgba(255,255,255,0.15);
        }
        .smooth-top-button:hover::after {
            opacity: 1;
            transform: translateX(0);
        }

        @media (max-width: 1200px) {
            .kpi-grid { grid-template-columns: repeat(2, minmax(170px, 1fr)); }
            .hero-top { flex-direction: column; align-items: flex-start; }
            .hero-meta { text-align: left; min-width: 0; }
        }
        </style>
        """,
        unsafe_allow_html=True,
    )

# ============================================================
# DATA HELPERS
# ============================================================
def find_column(df: pd.DataFrame, candidates: List[str]) -> Optional[str]:
    for c in candidates:
        if c in df.columns:
            return c
    return None


def normalize_dataframe(df: pd.DataFrame) -> pd.DataFrame:
    df = df.copy()
    df.columns = [str(c).strip().lower().replace(" ", "_") for c in df.columns]
    return df


def load_data(uploaded_file) -> pd.DataFrame:
    try:
        df = pd.read_csv(uploaded_file)
    except Exception:
        df = pd.read_csv(uploaded_file, encoding="utf-8", encoding_errors="replace")
    return normalize_dataframe(df)


def get_story_points_column(df): return find_column(df, ["story_points", "storypoints", "points", "estimate", "estimation", "story_point"])
def get_status_column(df): return find_column(df, ["status", "ticket_status", "issue_status"])
def get_type_column(df): return find_column(df, ["type", "issue_type", "issuetype"])
def get_priority_column(df): return find_column(df, ["priority", "issue_priority", "ticket_priority", "prio"])
def get_assignee_column(df): return find_column(df, ["assignee", "owner", "assigned_to", "resource"])
def get_sprint_column(df): return find_column(df, ["sprint", "sprint_name", "sprint_id", "iteration"])
def get_blocked_column(df): return find_column(df, ["blocked", "is_blocked", "blocked_flag", "flagged"])
def get_created_column(df): return find_column(df, ["created_date", "created", "start_date"])
def get_updated_column(df): return find_column(df, ["updated_date", "updated", "last_updated", "modified_date", "modified"])
def get_resolved_column(df): return find_column(df, ["resolved_date", "resolved", "closed_date", "end_date"])
def get_dependency_column(df): return find_column(df, ["dependency", "dependencies", "depends_on", "linked_issue", "blocked_by"])
def get_key_column(df): return find_column(df, ["key", "issue_key", "id", "ticket_id"])


def parse_dates(df: pd.DataFrame) -> pd.DataFrame:
    df = df.copy()
    for c in df.columns:
        if "date" in c or c in ["created", "resolved", "updated", "modified", "start_date", "end_date"]:
            df[c] = pd.to_datetime(df[c], errors="coerce")
    return df


def derive_time_dimensions(df: pd.DataFrame) -> pd.DataFrame:
    df = df.copy()
    date_col = get_created_column(df) or get_updated_column(df)
    if date_col:
        df[date_col] = pd.to_datetime(df[date_col], errors="coerce")
        df["filter_date"] = df[date_col]
    else:
        df["filter_date"] = pd.NaT
    df["month"] = df["filter_date"].dt.strftime("%Y-%m").fillna("Unknown")
    df["week"] = df["filter_date"].dt.strftime("%Y-W%V").fillna("Unknown")
    return df


def proper_case_value(v: Any) -> str:
    if pd.isna(v):
        return "Unknown"
    s = str(v).strip()
    if not s:
        return "Unknown"
    # Preserve sprint names like Sprint 18, IDs like D365-111
    if any(ch.isdigit() for ch in s) or "-" in s:
        return s
    return s.title()


def prepare_model(df: pd.DataFrame) -> pd.DataFrame:
    df = parse_dates(df)
    df = derive_time_dimensions(df)
    df = df.copy()

    sp_col = get_story_points_column(df)
    status_col = get_status_column(df)
    type_col = get_type_column(df)
    priority_col = get_priority_column(df)
    assignee_col = get_assignee_column(df)
    sprint_col = get_sprint_column(df)
    blocked_col = get_blocked_column(df)
    created_col = get_created_column(df)
    resolved_col = get_resolved_column(df)
    dep_col = get_dependency_column(df)
    key_col = get_key_column(df)

    df["story_points"] = pd.to_numeric(df[sp_col], errors="coerce").fillna(0) if sp_col else 0
    df["status"] = df[status_col].map(proper_case_value) if status_col else "Unknown"
    df["issue_type"] = df[type_col].map(proper_case_value) if type_col else "Unknown"
    df["priority"] = df[priority_col].map(proper_case_value) if priority_col else "Unknown"
    df["assignee"] = df[assignee_col].map(proper_case_value) if assignee_col else "Unassigned"
    df["sprint"] = df[sprint_col].map(proper_case_value) if sprint_col else "Unassigned"
    df["ticket_key"] = df[key_col].astype(str).fillna("") if key_col else df.index.astype(str)

    status_lower = df["status"].str.lower()
    done_mask = status_lower.isin(["done", "closed", "resolved", "completed"])

    blocked_mask = status_lower.eq("blocked")
    if blocked_col:
        b = df[blocked_col].astype(str).str.lower().fillna("")
        blocked_mask = blocked_mask | b.isin(["true", "yes", "blocked", "1"]) | b.str.contains("blocked", na=False)
    df["blocked"] = blocked_mask

    if created_col:
        df["created_date"] = pd.to_datetime(df[created_col], errors="coerce")
    else:
        df["created_date"] = pd.NaT
    if resolved_col:
        df["resolved_date"] = pd.to_datetime(df[resolved_col], errors="coerce")
    else:
        df["resolved_date"] = pd.NaT

    today = pd.Timestamp(datetime.today().date())
    df["days_open"] = np.where(
        df["created_date"].notna() & df["resolved_date"].notna(),
        (df["resolved_date"] - df["created_date"]).dt.days,
        np.where(df["created_date"].notna(), (today - df["created_date"]).dt.days, 0),
    )
    df["days_open"] = pd.Series(df["days_open"]).clip(lower=0).fillna(0).astype(int)

    if dep_col:
        dep_text = df[dep_col].fillna("").astype(str).str.strip()
        df["has_dependency"] = dep_text.ne("") & dep_text.str.lower().ne("nan")
        df["dependency_text"] = dep_text
    else:
        df["has_dependency"] = False
        df["dependency_text"] = ""

    df["risk_category"] = np.where(df["blocked"], "Blocked", np.where(df["days_open"] > 10, "Aging", "Normal"))
    df["is_done"] = done_mask
    df["is_defect"] = df["issue_type"].str.lower().eq("bug")
    df["is_high_risk"] = df["blocked"] | (df["days_open"] > 10) | df["priority"].str.lower().isin(["critical", "high"])
    return df


def compute_agile_metrics(df: pd.DataFrame) -> Dict[str, Any]:
    df = prepare_model(df)
    if df.empty:
        sprint_summary = pd.DataFrame()
        resource_workload = pd.DataFrame()
        defect_trend = pd.DataFrame()
    else:
        sprint_summary = (
            df.groupby("sprint", dropna=False)
            .agg(
                planned_story_points=("story_points", "sum"),
                completed_story_points=("story_points", lambda x: x[df.loc[x.index, "is_done"]].sum()),
                spillover_story_points=("story_points", lambda x: x[~df.loc[x.index, "is_done"]].sum()),
                defect_count=("is_defect", "sum"),
                blocked_tickets=("blocked", "sum"),
                aging_tickets=("days_open", lambda x: (x > 10).sum()),
                ticket_count=("ticket_key", "count"),
            )
            .reset_index()
        )
        sprint_summary["planned_vs_completed_pct"] = np.where(
            sprint_summary["planned_story_points"] > 0,
            sprint_summary["completed_story_points"] / sprint_summary["planned_story_points"] * 100,
            0,
        )
        sprint_summary["spillover_pct"] = np.where(
            sprint_summary["planned_story_points"] > 0,
            sprint_summary["spillover_story_points"] / sprint_summary["planned_story_points"] * 100,
            0,
        )
        sprint_summary["velocity"] = sprint_summary["completed_story_points"]
        sprint_summary["predictability"] = sprint_summary["planned_vs_completed_pct"]

        resource_workload = (
            df.groupby("assignee", dropna=False)
            .agg(
                workload_story_points=("story_points", "sum"),
                open_tickets=("is_done", lambda x: (~x).sum()),
                ticket_count=("ticket_key", "count"),
                blocked_tickets=("blocked", "sum"),
            )
            .reset_index()
            .sort_values("workload_story_points", ascending=False)
        )
        defect_trend = (
            df[df["is_defect"]]
            .groupby("sprint", dropna=False)
            .size()
            .reset_index(name="defect_count")
        )

    total_planned = float(df["story_points"].sum()) if not df.empty else 0
    total_completed = float(df.loc[df["is_done"], "story_points"].sum()) if not df.empty else 0
    total_spillover = total_planned - total_completed
    total_defects = int(df["is_defect"].sum()) if not df.empty else 0
    total_blocked = int(df["blocked"].sum()) if not df.empty else 0
    total_aging = int((df["days_open"] > 10).sum()) if not df.empty else 0
    total_dependencies = int(df["has_dependency"].sum() + df["blocked"].sum()) if not df.empty else 0
    ticket_count = len(df)

    planned_vs_completed_pct = total_completed / total_planned * 100 if total_planned > 0 else 0
    spillover_pct = total_spillover / total_planned * 100 if total_planned > 0 else 0
    avg_velocity = float(sprint_summary["velocity"].mean()) if not sprint_summary.empty else 0
    avg_predictability = float(sprint_summary["predictability"].mean()) if not sprint_summary.empty else 0

    defect_rate = total_defects / max(1, ticket_count) * 100
    blocked_rate = total_blocked / max(1, ticket_count) * 100
    aging_rate = total_aging / max(1, ticket_count) * 100

    health_score = 100
    health_score -= min(45, spillover_pct)
    health_score -= min(20, defect_rate)
    health_score -= min(20, blocked_rate)
    health_score -= min(15, aging_rate)
    health_score = int(max(0, min(100, round(health_score))))
    rag_status = "Green" if health_score >= 75 else "Amber" if health_score >= 50 else "Red"

    return {
        "df": df,
        "sprint_summary": sprint_summary,
        "resource_workload": resource_workload,
        "defect_trend": defect_trend,
        "agg_metrics": {
            "total_tickets": ticket_count,
            "total_sprints": int(df["sprint"].nunique()) if not df.empty else 0,
            "total_planned_story_points": total_planned,
            "total_completed_story_points": total_completed,
            "total_spillover_story_points": total_spillover,
            "planned_vs_completed_pct": planned_vs_completed_pct,
            "spillover_pct": spillover_pct,
            "defect_count": total_defects,
            "defect_rate": defect_rate,
            "blocked_count": total_blocked,
            "aging_tickets": total_aging,
            "dependency_count": total_dependencies,
            "health_score": health_score,
            "rag_status": rag_status,
            "average_velocity": avg_velocity,
            "average_predictability": avg_predictability,
            "resources_count": int(df["assignee"].nunique()) if not df.empty else 0,
            "high_risk_count": int(df["is_high_risk"].sum()) if not df.empty else 0,
        },
    }

# ============================================================
# FILTER STATE AND CROSS FILTERING
# ============================================================
FILTER_DEFAULTS = {
    "f_sprint": "All",
    "f_status": "All",
    "f_priority": "All",
    "f_assignee": "All",
    "f_month": "All",
    "f_week": "All",
    "f_risk_category": "All",
}


def init_filter_state():
    for k, v in FILTER_DEFAULTS.items():
        if k not in st.session_state:
            st.session_state[k] = v


def clear_filters():
    for k, v in FILTER_DEFAULTS.items():
        st.session_state[k] = v


# ============================================================
# JIRA TEMPLATE DOWNLOAD
# ============================================================
JIRA_TEMPLATE_COLUMNS = [
    "Key",
    "Sprint",
    "Type",
    "Status",
    "Assignee",
    "Priority",
    "Story Points",
    "Days Open",
    "Dependency",
    "Created",
    "Updated",
    "Resolved",
    "Blocked",
    "Labels",
]


def build_jira_template_dataframe() -> pd.DataFrame:
    """Create a Jira CSV template with example rows for Delivery Managers.

    The application accepts these columns and also supports close variants such as
    issue_type, story_points, assigned_to, blocked_by, created_date, etc. This
    template uses user-friendly Jira export column names.
    """
    rows = [
        {
            "Key": "D365-101",
            "Sprint": "Sprint 18",
            "Type": "Story",
            "Status": "Done",
            "Assignee": "John",
            "Priority": "High",
            "Story Points": 8,
            "Days Open": 4,
            "Dependency": "",
            "Created": "2026-06-01",
            "Updated": "2026-06-05",
            "Resolved": "2026-06-05",
            "Blocked": "No",
            "Labels": "frontend,release",
        },
        {
            "Key": "D365-102",
            "Sprint": "Sprint 18",
            "Type": "Task",
            "Status": "Blocked",
            "Assignee": "Smith",
            "Priority": "High",
            "Story Points": 5,
            "Days Open": 12,
            "Dependency": "D365-160",
            "Created": "2026-06-01",
            "Updated": "2026-06-13",
            "Resolved": "",
            "Blocked": "Yes",
            "Labels": "dependency,environment",
        },
        {
            "Key": "D365-103",
            "Sprint": "Sprint 18",
            "Type": "Bug",
            "Status": "Open",
            "Assignee": "Ravi",
            "Priority": "Medium",
            "Story Points": 3,
            "Days Open": 8,
            "Dependency": "",
            "Created": "2026-06-03",
            "Updated": "2026-06-11",
            "Resolved": "",
            "Blocked": "No",
            "Labels": "quality,defect",
        },
        {
            "Key": "D365-104",
            "Sprint": "Sprint 19",
            "Type": "Story",
            "Status": "In Progress",
            "Assignee": "Maria",
            "Priority": "Medium",
            "Story Points": 13,
            "Days Open": 6,
            "Dependency": "D365-210",
            "Created": "2026-06-07",
            "Updated": "2026-06-13",
            "Resolved": "",
            "Blocked": "No",
            "Labels": "integration",
        },
        {
            "Key": "D365-105",
            "Sprint": "Sprint 19",
            "Type": "Task",
            "Status": "Done",
            "Assignee": "Lee",
            "Priority": "Low",
            "Story Points": 2,
            "Days Open": 3,
            "Dependency": "",
            "Created": "2026-06-08",
            "Updated": "2026-06-11",
            "Resolved": "2026-06-11",
            "Blocked": "No",
            "Labels": "documentation",
        },
        {
            "Key": "D365-106",
            "Sprint": "Sprint 19",
            "Type": "Bug",
            "Status": "Blocked",
            "Assignee": "John",
            "Priority": "Critical",
            "Story Points": 8,
            "Days Open": 15,
            "Dependency": "Vendor API",
            "Created": "2026-06-02",
            "Updated": "2026-06-17",
            "Resolved": "",
            "Blocked": "Yes",
            "Labels": "vendor,critical,defect",
        },
    ]
    return pd.DataFrame(rows, columns=JIRA_TEMPLATE_COLUMNS)


def build_jira_template_csv_bytes() -> bytes:
    buffer = io.StringIO()
    build_jira_template_dataframe().to_csv(buffer, index=False)
    return buffer.getvalue().encode("utf-8")


def build_blank_jira_template_csv_bytes() -> bytes:
    buffer = io.StringIO()
    pd.DataFrame(columns=JIRA_TEMPLATE_COLUMNS).to_csv(buffer, index=False)
    return buffer.getvalue().encode("utf-8")


def get_options(series: pd.Series) -> List[str]:
    values = series.fillna("Unknown").astype(str).str.strip().replace({"": "Unknown"})
    return ["All"] + sorted(values.unique().tolist())


def apply_filters(df: pd.DataFrame) -> pd.DataFrame:
    filtered = df.copy()
    mapping = {
        "f_sprint": "sprint",
        "f_status": "status",
        "f_priority": "priority",
        "f_assignee": "assignee",
        "f_month": "month",
        "f_week": "week",
        "f_risk_category": "risk_category",
    }
    for state_key, col in mapping.items():
        selected = st.session_state.get(state_key, "All")
        if selected != "All" and col in filtered.columns:
            filtered = filtered[filtered[col].astype(str) == str(selected)]
    return filtered


def selected_filters_text() -> str:
    labels = []
    names = {
        "f_sprint": "Sprint",
        "f_status": "Status",
        "f_priority": "Priority",
        "f_assignee": "Assignee",
        "f_month": "Month",
        "f_week": "Week",
        "f_risk_category": "Risk",
    }
    for k, name in names.items():
        v = st.session_state.get(k, "All")
        if v != "All":
            labels.append(f"{name} = {v}")
    return " | ".join(labels) if labels else "All data selected"


def get_selected_point_value(event: Any, fallback_key: str = "x") -> Optional[str]:
    """Extract selected value from st.plotly_chart selection event."""
    try:
        if not event:
            return None
        selection = event.get("selection", None) if isinstance(event, dict) else getattr(event, "selection", None)
        if not selection:
            return None
        points = selection.get("points", None) if isinstance(selection, dict) else getattr(selection, "points", None)
        if not points:
            return None
        p0 = points[0]
        if isinstance(p0, dict):
            if "customdata" in p0 and p0["customdata"]:
                cd = p0["customdata"]
                if isinstance(cd, list) and len(cd) >= 2:
                    return str(cd[1])
                if isinstance(cd, str):
                    return cd
            return str(p0.get(fallback_key)) if p0.get(fallback_key) is not None else None
        return None
    except Exception:
        return None


def apply_chart_selection(event: Any, state_key: str, valid_options: List[str]) -> None:
    val = get_selected_point_value(event)
    if val and val in valid_options and st.session_state.get(state_key) != val:
        st.session_state[state_key] = val
        st.rerun()

# ============================================================
# VISUAL HELPERS
# ============================================================
def chart_layout(
    fig: go.Figure,
    title: str,
    x_title: str,
    y_title: str,
    height: int = CHART_HEIGHT,
    show_legend: bool = False,
) -> go.Figure:
    """Apply consistent dark executive chart styling.

    Legends are hidden by default to keep the client view clean.
    Only the Planned vs Completed chart uses a legend, positioned at the
    extreme top-right inside the reserved top margin.
    """
    top_margin = 118 if show_legend else 82
    fig.update_layout(
        title=dict(text=title, font=dict(color=TEXT, size=18), x=0.02, y=0.98, xanchor="left", yanchor="top"),
        plot_bgcolor=DARK_BG,
        paper_bgcolor=DARK_BG,
        font=dict(color=TEXT, size=12),
        height=height,
        margin=dict(l=72, r=42, t=top_margin, b=82),
        showlegend=show_legend,
        legend=dict(
            orientation="h",
            x=0.995,
            y=1.095,
            xanchor="right",
            yanchor="top",
            bgcolor="rgba(0,0,0,0)",
            font=dict(color=TEXT, size=11),
            title=dict(text="", font=dict(color=TEXT, size=11)),
            tracegroupgap=8,
            itemwidth=30,
        ),
        hoverlabel=dict(bgcolor="#FFFFFF", font_size=12, font_color="#0B1220"),
        bargap=0.32,
    )
    fig.update_xaxes(
        title_text=x_title,
        title_font=dict(color=TEXT, size=13),
        tickfont=dict(color=TEXT, size=11),
        gridcolor="rgba(255,255,255,0.06)",
        zerolinecolor=GRID,
        automargin=True,
    )
    fig.update_yaxes(
        title_text=y_title,
        title_font=dict(color=TEXT, size=13),
        tickfont=dict(color=TEXT, size=11),
        gridcolor=GRID,
        zerolinecolor=GRID,
        automargin=True,
    )
    return fig


def capitalize_column_name(col: str) -> str:
    """Convert technical column names into client-friendly display labels."""
    special = {
        "id": "ID",
        "sp": "SP",
        "kpi": "KPI",
        "rag": "RAG",
        "url": "URL",
        "ai": "AI",
        "jira": "Jira",
    }
    parts = str(col).replace("_", " ").replace("-", " ").split()
    return " ".join(special.get(part.lower(), part.capitalize()) for part in parts)


def display_table(df: pd.DataFrame) -> pd.DataFrame:
    out = df.copy()
    out.columns = [capitalize_column_name(c) for c in out.columns]
    return out


def render_chart(fig: go.Figure, key: str):
    st.plotly_chart(
        fig,
        use_container_width=True,
        key=key,
        config={"displayModeBar": False, "responsive": True},
    )


def add_customdata(fig: go.Figure, filter_name: str):
    for trace in fig.data:
        x_vals = list(trace.x) if hasattr(trace, "x") and trace.x is not None else []
        trace.customdata = [[filter_name, str(x)] for x in x_vals]


def kpi_card(title: str, value: str, subtitle: str, accent: str, tooltip: str) -> str:
    return f"""
    <div class='kpi-card' style='--accent:{accent};' title='{tooltip}'>
      <div class='kpi-title'>{title}</div>
      <div class='kpi-value'>{value}</div>
      <div class='kpi-subtitle'>{subtitle}</div>
    </div>
    """


def render_kpi_grid(agg: Dict[str, Any]):
    kpis = [
        ("Total Tickets", f"{agg['total_tickets']}", "Current filtered scope", BLUE, "Total Jira items included in the current filtered view."),
        ("Story Points", f"{int(agg['total_planned_story_points'])}", "Total planned effort", CYAN, "Total story points planned for the selected scope."),
        ("Completed SP", f"{int(agg['total_completed_story_points'])}", "Delivered effort", GREEN, "Story points completed based on Done/Closed/Resolved/Completed status."),
        ("Spillover SP", f"{int(agg['total_spillover_story_points'])}", f"{agg['spillover_pct']:.1f}% remaining", AMBER, "Story points not completed in the selected view."),
        ("RAG", agg["rag_status"], "Health indicator", RED if agg["rag_status"] == "Red" else AMBER if agg["rag_status"] == "Amber" else GREEN, "Overall delivery RAG derived from spillover, defects, blockers and aging tickets."),
        ("Avg Velocity", f"{agg['average_velocity']:.1f}", "Story points", CYAN, "Average completed story points across selected sprints."),
        ("Predictability", f"{agg['average_predictability']:.1f}%", "Average completion", BLUE, "Average planned versus completed percentage across selected sprints."),
        ("Blocked", f"{agg['blocked_count']}", "Current blockers", RED, "Tickets blocked by status or blocked flag."),
        ("Aging", f"{agg['aging_tickets']}", ">10 days open", AMBER, "Tickets open for more than 10 days."),
        ("Defects", f"{agg['defect_count']}", "Bug count", PURPLE, "Total bug items in the current filtered view."),
    ]
    for row_start in range(0, len(kpis), 5):
        cols = st.columns(5, gap="medium")
        for col, item in zip(cols, kpis[row_start:row_start + 5]):
            title, value, subtitle, accent, tooltip = item
            with col:
                st.markdown(kpi_card(title, value, subtitle, accent, tooltip), unsafe_allow_html=True)

def empty_chart(title: str, message: str) -> go.Figure:
    fig = go.Figure()
    fig.add_annotation(text=message, x=0.5, y=0.5, showarrow=False, font=dict(color=TEXT, size=16))
    return chart_layout(fig, title, "", "")


def make_status_chart(df: pd.DataFrame) -> go.Figure:
    if df.empty:
        return empty_chart("Status Distribution", "No data available")
    data = df["status"].value_counts().reset_index()
    data.columns = ["Status", "Count"]
    fig = px.bar(data, x="Status", y="Count", text="Count", color="Status", color_discrete_map=STATUS_COLOR_MAP)
    fig.update_traces(textposition="outside", marker_line_width=0, hovertemplate="Status: %{x}<br>Count: %{y}<extra></extra>")
    add_customdata(fig, "status")
    return chart_layout(fig, "Status Distribution", "Status", "Count")


def make_risk_chart(df: pd.DataFrame) -> go.Figure:
    if df.empty:
        return empty_chart("Risk Category Count", "No data available")
    data = df["risk_category"].value_counts().reset_index()
    data.columns = ["Risk Category", "Count"]
    fig = px.bar(data, x="Risk Category", y="Count", text="Count", color="Risk Category", color_discrete_map=RISK_COLOR_MAP)
    fig.update_traces(textposition="outside", marker_line_width=0, hovertemplate="Risk Category: %{x}<br>Count: %{y}<extra></extra>")
    add_customdata(fig, "risk_category")
    return chart_layout(fig, "Risk Category Count", "Risk Category", "Count")


def make_priority_chart(df: pd.DataFrame) -> go.Figure:
    if df.empty:
        return empty_chart("Priority Distribution", "No data available")
    data = df["priority"].value_counts().reset_index()
    data.columns = ["Priority", "Count"]
    fig = px.bar(data, x="Priority", y="Count", text="Count", color="Priority", color_discrete_map=PRIORITY_COLOR_MAP)
    fig.update_traces(textposition="outside", marker_line_width=0, hovertemplate="Priority: %{x}<br>Count: %{y}<extra></extra>")
    add_customdata(fig, "priority")
    return chart_layout(fig, "Priority Distribution", "Priority", "Count")


def make_velocity_chart(sprint_summary: pd.DataFrame) -> go.Figure:
    if sprint_summary.empty:
        return empty_chart("Sprint Velocity By Sprint", "No sprint data available")
    fig = px.bar(sprint_summary, x="sprint", y="velocity", text="velocity")
    fig.update_traces(marker_color=PRIMARY, texttemplate="%{text:.0f}", textposition="outside", hovertemplate="Sprint: %{x}<br>Velocity: %{y:.0f} SP<extra></extra>")
    add_customdata(fig, "sprint")
    return chart_layout(fig, "Sprint Velocity By Sprint", "Sprint", "Story Points")


def make_defect_trend(defect_trend: pd.DataFrame, sprint_summary: pd.DataFrame) -> go.Figure:
    if sprint_summary.empty:
        return empty_chart("Defect Trend", "No sprint data available")
    base = pd.DataFrame({"sprint": sprint_summary["sprint"].unique()})
    if defect_trend.empty:
        base["defect_count"] = 0
    else:
        base = base.merge(defect_trend, on="sprint", how="left").fillna({"defect_count": 0})
    fig = px.line(base, x="sprint", y="defect_count", markers=True, text="defect_count")
    fig.update_traces(line=dict(color=RED, width=3), marker=dict(size=8, color=RED), textposition="top center", hovertemplate="Sprint: %{x}<br>Defects: %{y:.0f}<extra></extra>")
    add_customdata(fig, "sprint")
    return chart_layout(fig, "Defect Trend", "Sprint", "Defect Count")


def make_planned_completed_chart(sprint_summary: pd.DataFrame) -> go.Figure:
    if sprint_summary.empty:
        return empty_chart("Planned Vs Completed Story Points", "No sprint data available")
    plot_df = sprint_summary[["sprint", "planned_story_points", "completed_story_points"]].rename(
        columns={"planned_story_points": "Planned Story Points", "completed_story_points": "Completed Story Points"}
    )
    long_df = plot_df.melt(id_vars="sprint", var_name="Metric", value_name="Story Points")
    fig = px.bar(
        long_df,
        x="sprint",
        y="Story Points",
        color="Metric",
        text="Story Points",
        barmode="group",
        color_discrete_map={"Planned Story Points": PRIMARY, "Completed Story Points": GREEN},
    )
    fig.update_traces(texttemplate="%{text:.0f}", textposition="outside", hovertemplate="Sprint: %{x}<br>%{fullData.name}: %{y:.0f}<extra></extra>")
    add_customdata(fig, "sprint")
    return chart_layout(fig, "Planned Vs Completed Story Points", "Sprint", "Story Points", show_legend=True)


def make_dependency_chart(df: pd.DataFrame) -> go.Figure:
    blocked_by_owner = df[df["blocked"]].groupby("assignee").size().reset_index(name="Blocked Dependencies").sort_values("Blocked Dependencies", ascending=False)
    if blocked_by_owner.empty:
        return empty_chart("Dependency Count By Assignee", "No blocked dependencies in selected view")
    fig = px.bar(blocked_by_owner.head(10), x="assignee", y="Blocked Dependencies", text="Blocked Dependencies")
    fig.update_traces(marker_color=PRIMARY, textposition="outside", hovertemplate="Assignee: %{x}<br>Blocked Dependencies: %{y}<extra></extra>")
    add_customdata(fig, "assignee")
    return chart_layout(fig, "Dependency Count By Assignee", "Assignee", "Blocked Dependencies")


def make_resource_chart(resource_workload: pd.DataFrame) -> go.Figure:
    if resource_workload.empty:
        return empty_chart("Resource Workload", "No resource data available")
    fig = px.bar(resource_workload.head(10), x="assignee", y="workload_story_points", text="workload_story_points")
    fig.update_traces(marker_color=PRIMARY, texttemplate="%{text:.0f}", textposition="outside", hovertemplate="Assignee: %{x}<br>Story Points: %{y:.0f}<extra></extra>")
    add_customdata(fig, "assignee")
    return chart_layout(fig, "Resource Workload", "Assignee", "Story Points")


def render_quality_kpi(agg: Dict[str, Any]):
    defect_rate = agg["defect_rate"]
    if defect_rate <= 10:
        status = "Healthy"
        color = GREEN
    elif defect_rate <= 20:
        status = "Needs Attention"
        color = AMBER
    else:
        status = "Escalate"
        color = RED
    html = f"""
    <div class='quality-panel' title='Quality KPI measures defect percentage from the selected Jira scope.'>
      <div class='quality-eyebrow'>Quality KPI</div>
      <div class='quality-title'>Defect Rate</div>
      <div class='quality-value'>{defect_rate:.1f}%</div>
      <div class='quality-status' style='background:{color}; color:white;'>{status}</div>
      <div class='quality-desc'>
        {agg['defect_count']} defects identified from {agg['total_tickets']} Jira items in the current filtered view. This indicates current rework pressure and quality attention required for the sprint discussion.
      </div>
      <div class='quality-mini-grid'>
        <div class='quality-mini-card'><div class='quality-mini-value'>{agg['defect_count']}</div><div class='quality-mini-label'>Defects</div></div>
        <div class='quality-mini-card'><div class='quality-mini-value'>{agg['total_tickets']}</div><div class='quality-mini-label'>Total items</div></div>
        <div class='quality-mini-card'><div class='quality-mini-value'>≤ 10%</div><div class='quality-mini-label'>Healthy target</div></div>
        <div class='quality-mini-card'><div class='quality-mini-value'>10% - 20%</div><div class='quality-mini-label'>Monitor range</div></div>
        <div class='quality-mini-card'><div class='quality-mini-value'>&gt; 20%</div><div class='quality-mini-label'>Escalate range</div></div>
        <div class='quality-mini-card'><div class='quality-mini-value'>{status}</div><div class='quality-mini-label'>Current status</div></div>
      </div>
    </div>
    """
    st.markdown(html, unsafe_allow_html=True)

# ============================================================
# PPT GENERATOR
# ============================================================
def add_slide_bg(slide, color_hex=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(color_hex.replace("#", ""))


def add_textbox(slide, text, left, top, width, height, font_size=18, color="FFFFFF", bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(font_size)
    run.font.color.rgb = RGBColor.from_string(color.replace("#", ""))
    run.font.bold = bold
    if align:
        p.alignment = align
    return box


def add_ppt_kpi(slide, title, value, subtitle, x, y, w, h, accent):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(CARD_BG.replace("#", ""))
    shape.line.color.rgb = RGBColor.from_string("2D4A6E")
    accent_shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(0.07), Inches(h))
    accent_shape.fill.solid()
    accent_shape.fill.fore_color.rgb = RGBColor.from_string(accent.replace("#", ""))
    accent_shape.line.fill.background()
    add_textbox(slide, title, x + 0.15, y + 0.12, w - 0.25, 0.25, 9, "D7E3F4", True)
    add_textbox(slide, value, x + 0.15, y + 0.45, w - 0.25, 0.45, 22, "FFFFFF", True)
    add_textbox(slide, subtitle, x + 0.15, y + 0.92, w - 0.25, 0.25, 8, "D7E3F4", False)



def generate_client_engagement_ppt(df: pd.DataFrame) -> bytes:
    """Generate a client engagement KPI pack using one consistent template.

    Slide sequence:
    1 Landing Page
    2 Index / Agenda
    3 Delivery Governance
    4 Executive KPI
    5 Delivery Health & Status
    6 Sprint Trend
    7 Risks, Blockers & Aging Facts
    8 Dependency Facts
    9 Resource Workload
    10 Quality & Defect
    11 RAID / Action Tracker
    12 Client Decisions & Next Steps
    """
    metrics = compute_agile_metrics(df)
    agg = metrics["agg_metrics"]
    model_df = metrics["df"].copy()
    sprint_summary = metrics["sprint_summary"].copy()
    resource_workload = metrics["resource_workload"].copy()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    generated_on = datetime.now(ZoneInfo("Asia/Kolkata")).strftime("%d %b %Y, %I:%M %p IST")
    deck_name = "CLIENT ENGAGEMENT PACK | DELIVERY GOVERNANCE"

    # -------------------------
    # Consistent slide settings
    # -------------------------
    SLIDE_W = 13.333
    SLIDE_H = 7.5
    MARGIN_X = 0.55
    TITLE_Y = 0.32
    SUBTITLE_Y = 0.82
    CONTENT_TOP = 1.25
    FOOTER_Y = 7.08

    def hex_clean(color: str) -> str:
        return str(color).replace("#", "")

    def rag_color() -> str:
        return RED if agg["rag_status"] == "Red" else AMBER if agg["rag_status"] == "Amber" else GREEN

    def add_footer(slide, page_no: int):
        add_textbox(slide, deck_name, MARGIN_X, FOOTER_Y, 8.0, 0.22, 8, "D7E3F4", True)
        add_textbox(slide, str(page_no), 12.45, FOOTER_Y, 0.35, 0.22, 8, "D7E3F4", False, PP_ALIGN.RIGHT)

    def new_slide(title: str, subtitle: str, page_no: int):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_slide_bg(slide)
        add_textbox(slide, title.upper(), MARGIN_X, TITLE_Y, 12.1, 0.38, 22, "FFFFFF", True)
        if subtitle:
            add_textbox(slide, subtitle, MARGIN_X + 0.02, SUBTITLE_Y, 12.0, 0.25, 9.5, "D7E3F4", False)
        add_footer(slide, page_no)
        return slide

    def add_small_kpi(slide, title, value, subtitle, x, y, w=2.35, h=0.96, accent=BLUE):
        shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(hex_clean(CARD_BG))
        shape.line.color.rgb = RGBColor.from_string("2D4A6E")
        shape.line.width = Pt(0.8)
        bar = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(0.06), Inches(h))
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor.from_string(hex_clean(accent))
        bar.line.fill.background()
        add_textbox(slide, str(title).upper(), x + 0.16, y + 0.10, w - 0.23, 0.18, 7.5, "D7E3F4", True)
        value_size = 19 if len(str(value)) <= 7 else 15
        add_textbox(slide, str(value), x + 0.16, y + 0.36, w - 0.23, 0.30, value_size, "FFFFFF", True)
        add_textbox(slide, str(subtitle), x + 0.16, y + 0.69, w - 0.23, 0.18, 6.8, "D7E3F4", False)

    def add_section_label(slide, text, x, y, w=2.8):
        add_textbox(slide, text.upper(), x, y, w, 0.20, 8, hex_clean(CYAN), True)

    def add_text_lines(slide, title, lines, x, y, w, line_h=0.29, font_size=8.5):
        if title:
            add_textbox(slide, title, x, y, w, 0.25, 11, hex_clean(CYAN), True)
            yy = y + 0.36
        else:
            yy = y
        for line in lines:
            add_textbox(slide, str(line), x + 0.03, yy, w - 0.05, 0.24, font_size, "D7E3F4", False)
            yy += line_h

    def add_ppt_table(slide, df_table: pd.DataFrame, x, y, w, h, font_size=6.6, header_font_size=6.8, max_rows=8):
        if df_table is None or df_table.empty:
            add_textbox(slide, "No records available for the current filter context.", x, y, w, 0.3, 9, "D7E3F4", False)
            return
        table_df = df_table.head(max_rows).copy().fillna("")
        rows = len(table_df) + 1
        cols = len(table_df.columns)
        shape = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h))
        table = shape.table
        for c, col in enumerate(table_df.columns):
            cell = table.cell(0, c)
            cell.text = str(col)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor.from_string("102B46")
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.size = Pt(header_font_size)
                    run.font.bold = True
                    run.font.color.rgb = RGBColor.from_string("FFFFFF")
        for r in range(len(table_df)):
            for c, col in enumerate(table_df.columns):
                val = table_df.iloc[r, c]
                if isinstance(val, float):
                    text = f"{val:.1f}" if abs(val) < 1000 else f"{val:,.0f}"
                else:
                    text = str(val)
                cell = table.cell(r + 1, c)
                cell.text = text
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor.from_string("0B2138" if r % 2 == 0 else "0D2742")
                for p in cell.text_frame.paragraphs:
                    p.alignment = PP_ALIGN.LEFT
                    for run in p.runs:
                        run.font.size = Pt(font_size)
                        run.font.color.rgb = RGBColor.from_string("D7E3F4")

    def count_table(series: pd.Series, first_col: str) -> pd.DataFrame:
        out = series.fillna("Unknown").astype(str).value_counts().reset_index()
        out.columns = [first_col, "Count"]
        return out

    def risk_score_row(row) -> int:
        score = 0
        if bool(row.get("blocked", False)):
            score += 55
        priority = str(row.get("priority", "")).lower()
        if priority == "critical":
            score += 35
        elif priority == "high":
            score += 20
        if float(row.get("days_open", 0)) > 10:
            score += 25
        return int(min(100, score))

    if model_df.empty:
        model_df = pd.DataFrame(columns=["ticket_key", "sprint", "status", "priority", "assignee", "story_points", "days_open", "blocked", "dependency_text", "issue_type"])
    model_df["Risk Score"] = model_df.apply(risk_score_row, axis=1) if not model_df.empty else []

    completed_tickets = int(model_df["is_done"].sum()) if "is_done" in model_df else 0
    open_tickets = int(len(model_df) - completed_tickets)
    completion_rate = completed_tickets / max(1, len(model_df)) * 100
    planned_completed_sp_pct = agg["planned_vs_completed_pct"]
    quality_impact = min(20, agg["defect_rate"]) / 5
    avg_workload = agg["total_planned_story_points"] / max(1, agg["resources_count"])
    overloaded_count = 0
    if not resource_workload.empty and avg_workload:
        overloaded_count = int((resource_workload["workload_story_points"] >= 1.5 * avg_workload).sum())

    # -------------------------
    # Slide 1: Landing Page - Executive redesign
    # -------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)

    def add_cover_kpi(slide, title, value, subtitle, x, y, w, h, accent):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor.from_string(hex_clean(CARD_BG))
        card.line.color.rgb = RGBColor.from_string("365B85")
        card.line.width = Pt(1.1)
        bar = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(0.08), Inches(h))
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor.from_string(hex_clean(accent))
        bar.line.fill.background()
        add_textbox(slide, str(title).upper(), x + 0.24, y + 0.18, w - 0.38, 0.22, 8.4, "D7E3F4", True)
        value_font = 27 if len(str(value)) <= 6 else 22
        add_textbox(slide, str(value), x + 0.24, y + 0.50, w - 0.38, 0.48, value_font, "FFFFFF", True)
        add_textbox(slide, str(subtitle), x + 0.24, y + 1.02, w - 0.38, 0.22, 8.2, "D7E3F4", False)

    # Top executive identity
    add_textbox(slide, "AI AGENT - DELIVERY MANAGER", 0.70, 0.42, 12.0, 0.28, 12.5, hex_clean(CYAN), True, PP_ALIGN.CENTER)
    add_textbox(slide, "Client Engagement Meeting", 0.70, 1.02, 12.0, 0.36, 21, "D7E3F4", True, PP_ALIGN.CENTER)
    add_textbox(slide, "DELIVERY GOVERNANCE REPORT", 0.70, 1.48, 12.0, 0.62, 36, "FFFFFF", True, PP_ALIGN.CENTER)
    add_textbox(
        slide,
        "Executive Delivery Review | Sprint Health, Risks, Dependencies & Delivery Performance",
        1.35, 2.22, 10.6, 0.30, 12.5, "D7E3F4", False, PP_ALIGN.CENTER,
    )

    # Executive summary strip
    strip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.00), Inches(2.92), Inches(11.30), Inches(0.48))
    strip.fill.solid()
    strip.fill.fore_color.rgb = RGBColor.from_string("0B2138")
    strip.line.color.rgb = RGBColor.from_string("2D4A6E")
    summary_line = (
        f"Current Delivery Health: {agg['rag_status'].upper()}   |   "
        f"Completion: {completion_rate:.1f}%   |   "
        f"Blocked Items: {agg['blocked_count']}   |   "
        f"Dependencies: {agg['dependency_count']}"
    )
    add_textbox(slide, summary_line, 1.15, 3.05, 11.0, 0.18, 10.5, "FFFFFF", True, PP_ALIGN.CENTER)

    # Bigger, aligned executive KPI cards
    y_cards = 3.82
    add_cover_kpi(slide, "Health Score", agg["health_score"], agg["rag_status"], 0.80, y_cards, 2.75, 1.35, rag_color())
    add_cover_kpi(slide, "Completion", f"{completion_rate:.1f}%", "Ticket closure", 3.85, y_cards, 2.75, 1.35, GREEN if completion_rate >= 75 else AMBER)
    add_cover_kpi(slide, "Blockers", agg["blocked_count"], "Needs attention", 6.90, y_cards, 2.75, 1.35, RED)
    add_cover_kpi(slide, "Delivered SP", int(agg["total_completed_story_points"]), "Completed story points", 9.95, y_cards, 2.75, 1.35, CYAN)

    # Executive insight box
    insight = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.00), Inches(5.55), Inches(11.30), Inches(0.95))
    insight.fill.solid()
    insight.fill.fore_color.rgb = RGBColor.from_string("0B2138")
    insight.line.color.rgb = RGBColor.from_string("2D4A6E")
    add_textbox(slide, "EXECUTIVE INSIGHT", 1.22, 5.70, 2.2, 0.18, 8.5, hex_clean(CYAN), True)
    insight_text = (
        f"Delivery health is {agg['rag_status']} with {agg['blocked_count']} blocker(s), "
        f"{agg['aging_tickets']} aging item(s), and {agg['dependency_count']} dependency item(s). "
        f"Recommend weekly governance checkpoint until risks and RAID actions are closed."
    )
    add_textbox(slide, insight_text, 1.22, 5.96, 10.85, 0.28, 10.2, "D7E3F4", False)

    add_textbox(slide, f"Generated from current UI filters on {generated_on}", 1.00, 6.63, 11.30, 0.22, 8.5, "D7E3F4", False, PP_ALIGN.CENTER)
    add_footer(slide, 1)

    # -------------------------
    # Slide 2: Index / Agenda
    # -------------------------
    slide = new_slide("Index / Agenda", "Discussion flow for the client engagement meeting", 2)
    agenda = [
        "01  Delivery Governance",
        "02  Executive KPI",
        "03  Delivery Health & Status",
        "04  Sprint Trend",
        "05  Risks, Blockers & Aging Facts",
        "06  Dependency Facts",
        "07  Resource Workload",
        "08  Quality & Defect",
        "09  RAID / Action Tracker",
        "10  Client Decisions & Next Steps",
    ]
    left_items = agenda[:5]
    right_items = agenda[5:]
    def agenda_item(txt, x, y):
        box = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(5.7), Inches(0.52))
        box.fill.solid(); box.fill.fore_color.rgb = RGBColor.from_string(hex_clean(CARD_BG))
        box.line.color.rgb = RGBColor.from_string("2D4A6E")
        add_textbox(slide, txt, x+0.20, y+0.14, 5.2, 0.18, 12, "FFFFFF", True)
    y = 1.35
    for item in left_items:
        agenda_item(item, 0.85, y); y += 0.72
    y = 1.35
    for item in right_items:
        agenda_item(item, 6.80, y); y += 0.72

    # -------------------------
    # Slide 3: Delivery Governance
    # -------------------------
    slide = new_slide("Delivery Governance", "Fact-based governance view generated from the current Streamlit UI filters", 3)
    add_small_kpi(slide, "Current RAG", agg["rag_status"], f"Health Score {agg['health_score']}", 0.65, 1.28, 2.65, 1.02, rag_color())
    add_small_kpi(slide, "Governance Scope", agg["total_tickets"], "Jira items reviewed", 3.55, 1.28, 2.65, 1.02, BLUE)
    add_small_kpi(slide, "Client Focus", "Risks", "Blockers + dependencies", 6.45, 1.28, 2.65, 1.02, RED)
    add_small_kpi(slide, "Cadence", "Weekly", "Track to closure", 9.35, 1.28, 2.65, 1.02, CYAN)
    add_section_label(slide, "Governance Pillars", 0.70, 2.75)
    pillars = [
        ("Delivery Health", "Health score, RAG status, completion and open work."),
        ("Risk Control", "Blocked, aging and high-risk tickets requiring mitigation."),
        ("Dependency Management", "Dependency items, bottlenecks, owners and ETAs."),
        ("Resource Governance", "Workload concentration and capacity pressure."),
        ("Quality Governance", "Defects, defect rate and quality impact."),
        ("Action Tracking", "RAID actions, owners, dates and closure rhythm."),
    ]
    for i, (t, d) in enumerate(pillars):
        x = 0.75 + (i % 3) * 4.05
        y = 3.15 + (i // 3) * 1.15
        box = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(3.65), Inches(0.82))
        box.fill.solid(); box.fill.fore_color.rgb = RGBColor.from_string("0B2138")
        box.line.color.rgb = RGBColor.from_string("2D4A6E")
        add_textbox(slide, t, x+0.16, y+0.12, 3.3, 0.20, 10, "FFFFFF", True)
        add_textbox(slide, d, x+0.16, y+0.38, 3.3, 0.25, 7.5, "D7E3F4", False)

    # -------------------------
    # Slide 4: Executive KPI
    # -------------------------
    slide = new_slide("Executive KPI", "Datapoints exported from the current Streamlit filtered view", 4)
    kpis = [
        ("Total Tickets", agg["total_tickets"], "Current scope", BLUE),
        ("Story Points", int(agg["total_planned_story_points"]), "Total effort", CYAN),
        ("Completed SP", int(agg["total_completed_story_points"]), "Delivered", GREEN),
        ("Spillover SP", int(agg["total_spillover_story_points"]), f"{agg['spillover_pct']:.1f}%", AMBER),
        ("RAG", agg["rag_status"], "Health", rag_color()),
        ("Avg Velocity", f"{agg['average_velocity']:.1f}", "Story points", CYAN),
        ("Predictability", f"{agg['average_predictability']:.1f}%", "Avg sprint", BLUE),
        ("Blocked", agg["blocked_count"], "Current blockers", RED),
        ("Aging", agg["aging_tickets"], ">10 days", AMBER),
        ("Defects", agg["defect_count"], "Bug count", PURPLE),
    ]
    for i, k in enumerate(kpis):
        x = 0.65 + (i % 5) * 2.48
        y = 1.35 + (i // 5) * 1.18
        title, value, subtitle, accent = k
        add_small_kpi(slide, title, value, subtitle, x, y, w=2.18, h=0.96, accent=accent)
    facts = [
        f"Completion rate is {completion_rate:.1f}% based on {completed_tickets} completed tickets out of {agg['total_tickets']}.",
        f"Planned vs completed story points is {planned_completed_sp_pct:.1f}%.",
        f"Current delivery health is {agg['health_score']} / 100 with {agg['rag_status']} RAG status.",
        f"There are {agg['blocked_count']} blocked tickets and {agg['aging_tickets']} aging tickets requiring review.",
    ]
    add_text_lines(slide, "Key Facts", facts, 0.70, 4.25, 12.0, 0.30)

    # -------------------------
    # Slide 5: Delivery Health & Status
    # -------------------------
    slide = new_slide("Delivery Health & Status", "Ticket status, type distribution and delivery health indicators", 5)
    add_small_kpi(slide, "Health Score", agg["health_score"], f"RAG: {agg['rag_status']}", 0.65, 1.25, 2.75, 1.02, rag_color())
    add_small_kpi(slide, "Open Tickets", open_tickets, "Not completed", 3.65, 1.25, 2.75, 1.02, AMBER)
    add_small_kpi(slide, "Completed", completed_tickets, f"{completion_rate:.1f}%", 6.65, 1.25, 2.75, 1.02, GREEN)
    type_tbl = count_table(model_df["issue_type"], "Type") if "issue_type" in model_df else pd.DataFrame()
    status_tbl = count_table(model_df["status"], "Status") if "status" in model_df else pd.DataFrame()
    add_ppt_table(slide, type_tbl, 0.75, 3.00, 5.6, 2.45, font_size=7.5, header_font_size=7.5, max_rows=8)
    add_ppt_table(slide, status_tbl, 6.95, 3.00, 5.6, 2.45, font_size=7.5, header_font_size=7.5, max_rows=8)

    # -------------------------
    # Slide 6: Sprint Trend
    # -------------------------
    slide = new_slide("Sprint Trend", "Velocity, planned vs completed, predictability and spillover facts", 6)
    sprint_tbl = pd.DataFrame()
    if not sprint_summary.empty:
        sprint_tbl = sprint_summary[["sprint", "planned_story_points", "completed_story_points", "planned_vs_completed_pct", "spillover_story_points", "spillover_pct"]].copy()
        sprint_tbl.columns = ["Sprint", "Planned SP", "Completed SP", "Complete %", "Spillover SP", "Spillover %"]
        sprint_tbl["Complete %"] = sprint_tbl["Complete %"].map(lambda x: f"{x:.1f}%")
        sprint_tbl["Spillover %"] = sprint_tbl["Spillover %"].map(lambda x: f"{x:.1f}%")
    add_ppt_table(slide, sprint_tbl, 0.65, 1.35, 12.05, 4.85, font_size=7, header_font_size=7.2, max_rows=10)

    # -------------------------
    # Slide 7: Risks, Blockers & Aging Facts
    # -------------------------
    slide = new_slide("Risks, Blockers & Aging Facts", "Client discussion points from blocked, aging and high-priority items", 7)
    add_small_kpi(slide, "Blocked Tickets", agg["blocked_count"], "Status or blocked flag", 0.65, 1.25, 2.75, 1.02, RED)
    add_small_kpi(slide, "Aging Tickets", agg["aging_tickets"], ">10 days open", 3.65, 1.25, 2.75, 1.02, AMBER)
    add_small_kpi(slide, "Dependency Items", agg["dependency_count"], "Dependency/blocked facts", 6.65, 1.25, 2.75, 1.02, PURPLE)
    add_small_kpi(slide, "High Risk", agg["high_risk_count"], "Score >= 60", 9.65, 1.25, 2.75, 1.02, RED)
    risk_cols = ["ticket_key", "status", "priority", "assignee", "days_open", "blocked", "Risk Score"]
    risk_tbl = model_df.sort_values(["Risk Score", "days_open"], ascending=[False, False])[risk_cols].copy() if set(risk_cols).issubset(model_df.columns) else pd.DataFrame()
    if not risk_tbl.empty:
        risk_tbl.columns = ["Key", "Status", "Priority", "Owner", "Days Open", "Blocked", "Risk"]
    add_ppt_table(slide, risk_tbl, 0.65, 2.75, 12.05, 2.65, font_size=6.4, header_font_size=6.6, max_rows=8)
    add_text_lines(slide, "", [
        "Client focus: confirm recovery plan for all blocked and aging items.",
        "Priority focus: review high-priority items that are still open or blocked.",
        "Governance focus: capture owner, ETA and mitigation for each high-risk ticket.",
    ], 0.70, 5.65, 12.0, 0.24)

    # -------------------------
    # Slide 8: Dependency Facts
    # -------------------------
    slide = new_slide("Dependency Facts", "Dependency and bottleneck datapoints from the selected Jira view", 8)
    blocked_owners = int(model_df.loc[model_df["blocked"], "assignee"].nunique()) if "blocked" in model_df else 0
    blocked_sp = int(model_df.loc[model_df["blocked"], "story_points"].sum()) if "blocked" in model_df else 0
    add_small_kpi(slide, "Dependency Count", agg["dependency_count"], "Explicit dependency or blocked", 0.65, 1.25, 2.75, 1.02, PURPLE)
    add_small_kpi(slide, "Blocked Owners", blocked_owners, "Owners with blockers", 3.65, 1.25, 2.75, 1.02, RED)
    add_small_kpi(slide, "Blocked SP", blocked_sp, "Story points impacted", 6.65, 1.25, 2.75, 1.02, AMBER)
    dep_df = model_df[(model_df.get("has_dependency", False) == True) | (model_df.get("blocked", False) == True)].copy() if not model_df.empty else pd.DataFrame()
    if not dep_df.empty:
        dep_tbl = dep_df[["ticket_key", "dependency_text", "status", "assignee", "days_open"]].copy()
        dep_tbl.columns = ["Key", "Dependency", "Status", "Owner", "Days Open"]
    else:
        dep_tbl = pd.DataFrame()
    add_ppt_table(slide, dep_tbl, 0.65, 2.75, 12.05, 2.65, font_size=6.4, header_font_size=6.6, max_rows=8)
    add_text_lines(slide, "", [
        "Client decision required where external or cross-team dependencies are blocking sprint progress.",
        "All dependency items should have owner, ETA, mitigation and escalation path.",
    ], 0.70, 5.65, 12.0, 0.24)

    # -------------------------
    # Slide 9: Resource Workload
    # -------------------------
    slide = new_slide("Resource Workload", "Assignee workload, open tickets and capacity concentration", 9)
    add_small_kpi(slide, "Active Resources", agg["resources_count"], "Assignees", 0.65, 1.25, 2.75, 1.02, BLUE)
    add_small_kpi(slide, "Avg Workload", f"{avg_workload:.1f}", "SP per resource", 3.65, 1.25, 2.75, 1.02, CYAN)
    add_small_kpi(slide, "Overloaded", overloaded_count, ">=150% avg SP", 6.65, 1.25, 2.75, 1.02, RED if overloaded_count else GREEN)
    res_tbl = pd.DataFrame()
    if not resource_workload.empty:
        res_tbl = resource_workload[["assignee", "workload_story_points", "ticket_count", "open_tickets"]].copy()
        res_tbl.columns = ["Assignee", "Story Points", "Tickets", "Open"]
    add_ppt_table(slide, res_tbl, 0.65, 2.70, 9.80, 2.85, font_size=7, header_font_size=7.2, max_rows=8)

    # -------------------------
    # Slide 10: Quality & Defect
    # -------------------------
    slide = new_slide("Quality & Defect", "Defect count, defect rate and impact to sprint health", 10)
    add_small_kpi(slide, "Defects", agg["defect_count"], "Bug items", 0.65, 1.25, 2.75, 1.02, PURPLE)
    add_small_kpi(slide, "Defect Rate", f"{agg['defect_rate']:.1f}%", "Defects / tickets", 3.65, 1.25, 2.75, 1.02, AMBER if agg["defect_rate"] <= 20 else RED)
    add_small_kpi(slide, "Quality Impact", f"{quality_impact:.1f}", "Health deduction", 6.65, 1.25, 2.75, 1.02, RED if quality_impact > 3 else AMBER)
    bug_df = model_df[model_df["is_defect"]].copy() if "is_defect" in model_df else pd.DataFrame()
    if not bug_df.empty:
        bug_tbl = bug_df[["ticket_key", "status", "priority", "assignee", "days_open"]].copy()
        bug_tbl.columns = ["Key", "Status", "Priority", "Owner", "Days Open"]
    else:
        bug_tbl = pd.DataFrame()
    add_ppt_table(slide, bug_tbl, 0.65, 2.70, 9.80, 2.85, font_size=7, header_font_size=7.2, max_rows=8)

    # -------------------------
    # Slide 11: RAID / Action Tracker
    # -------------------------
    slide = new_slide("RAID / Action Tracker", "Fact-based actions for client engagement discussion", 11)
    raid_rows = []
    if not model_df.empty and "Risk Score" in model_df.columns:
        for _, row in model_df.sort_values(["Risk Score", "days_open"], ascending=[False, False]).head(5).iterrows():
            sev = "Critical" if int(row.get("Risk Score", 0)) >= 90 else "High" if int(row.get("Risk Score", 0)) >= 60 else "Medium"
            raid_rows.append({
                "Type": "Issue" if bool(row.get("blocked", False)) else "Risk",
                "ID": row.get("ticket_key", ""),
                "Fact": f"{row.get('status','')} | {row.get('priority','')} | {int(row.get('days_open',0))} days open",
                "Severity": sev,
                "Client Action / Mitigation": "Confirm owner, ETA and escalation path",
            })
    if not dep_df.empty:
        first_dep = dep_df.iloc[0]
        raid_rows.append({
            "Type": "Dependency",
            "ID": first_dep.get("dependency_text", first_dep.get("ticket_key", "")) or first_dep.get("ticket_key", ""),
            "Fact": "Linked dependency in current sprint",
            "Severity": "High",
            "Client Action / Mitigation": "Agree dependency owner and due date",
        })
    raid_tbl = pd.DataFrame(raid_rows)
    add_ppt_table(slide, raid_tbl, 0.55, 1.25, 12.25, 4.20, font_size=6.2, header_font_size=6.2, max_rows=7)
    add_text_lines(slide, "", [
        "Use this slide to agree owners and dates during the client engagement meeting.",
        "Update each item after the meeting and track weekly until closure.",
    ], 0.70, 5.75, 12.0, 0.26)

    # -------------------------
    # Slide 12: Client Decisions & Next Steps
    # -------------------------
    slide = new_slide("Client Decisions & Next Steps", "Recommended discussion points based on current facts and KPIs", 12)
    add_small_kpi(slide, "Current RAG", agg["rag_status"], f"Health Score {agg['health_score']}", 0.65, 1.25, 2.75, 1.02, rag_color())
    add_small_kpi(slide, "Client Focus", "Risks", "Blockers + dependencies", 3.65, 1.25, 2.75, 1.02, RED)
    add_small_kpi(slide, "Governance", "Weekly", "Track to closure", 6.65, 1.25, 2.75, 1.02, BLUE)
    next_steps = [
        f"Review {agg['blocked_count']} blocked tickets and agree recovery owner/date.",
        f"Review {agg['aging_tickets']} aging tickets above 10 days and agree closure plan.",
        f"Confirm mitigation for {agg['dependency_count']} dependency-risk item(s).",
        f"Review sprint completion at {planned_completed_sp_pct:.1f}% and spillover at {agg['spillover_pct']:.1f}%.",
        "Agree weekly checkpoint rhythm for risks, dependencies and RAID closure.",
    ]
    add_text_lines(slide, "Recommended Next Steps", next_steps, 0.70, 3.00, 12.0, 0.33)
    add_textbox(slide, "Thank You", 0.70, 6.10, 12.0, 0.45, 24, "FFFFFF", True, PP_ALIGN.CENTER)

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output.getvalue()

# ============================================================
# AI ASSISTANT
# ============================================================
def _fmt_num(value: Any, decimals: int = 1) -> str:
    try:
        f = float(value)
        if abs(f - int(f)) < 0.00001:
            return f"{int(f)}"
        return f"{f:.{decimals}f}"
    except Exception:
        return str(value)


def _ticket_list(df: pd.DataFrame, max_rows: int = 8) -> str:
    if df is None or df.empty:
        return "No matching Jira items found in the current filtered view."
    cols = [c for c in ["ticket_key", "sprint", "status", "priority", "assignee", "story_points", "days_open", "dependency_text"] if c in df.columns]
    rows = []
    for _, r in df[cols].head(max_rows).iterrows():
        key = r.get("ticket_key", "")
        sprint = r.get("sprint", "")
        status = r.get("status", "")
        priority = r.get("priority", "")
        owner = r.get("assignee", "")
        sp = _fmt_num(r.get("story_points", 0), 0)
        days = _fmt_num(r.get("days_open", 0), 0)
        dep = str(r.get("dependency_text", "")).strip()
        dep_text = f" | Dependency: {dep}" if dep else ""
        rows.append(f"- **{key}** | {sprint} | {status} | {priority} | Owner: {owner} | SP: {sp} | Days open: {days}{dep_text}")
    return "\n".join(rows)


def _governance_recommendations(agg: Dict[str, Any]) -> List[str]:
    recommendations = []
    if agg["rag_status"] == "Red":
        recommendations.append("Run a focused recovery checkpoint with client and delivery owners until the RAG moves out of Red.")
    elif agg["rag_status"] == "Amber":
        recommendations.append("Track the sprint through a weekly governance checkpoint and close key risks before the next review.")
    else:
        recommendations.append("Maintain current governance rhythm and continue monitoring leading indicators.")
    if agg["blocked_count"] > 0:
        recommendations.append(f"Prioritize the {agg['blocked_count']} blocked item(s); each should have owner, ETA, dependency owner and mitigation.")
    if agg["aging_tickets"] > 0:
        recommendations.append(f"Review {agg['aging_tickets']} aging item(s) above 10 days and agree closure or de-scope action.")
    if agg["dependency_count"] > 0:
        recommendations.append(f"Review {agg['dependency_count']} dependency-risk item(s) and confirm cross-team/client decision points.")
    if agg["defect_rate"] > 10:
        recommendations.append(f"Defect rate is {_fmt_num(agg['defect_rate'])}%; plan defect triage and root-cause review.")
    if agg["spillover_pct"] > 30:
        recommendations.append(f"Spillover is {_fmt_num(agg['spillover_pct'])}%; review commitment quality, scope change and capacity assumptions.")
    return recommendations[:6]


def answer_ai_question(question: str, metrics: Dict[str, Any]) -> str:
    """Rule-based governance assistant grounded only in the uploaded Jira data.

    It answers both data questions and governance/agile/sprint questions without
    calling external services, so it works in Streamlit Cloud without API keys.
    """
    q = question.lower().strip()
    df = metrics["df"].copy()
    agg = metrics["agg_metrics"]
    sprint_summary = metrics["sprint_summary"].copy()
    resource_workload = metrics["resource_workload"].copy()
    defect_trend = metrics["defect_trend"].copy()

    if df.empty:
        return "No Jira data is available for the current filter context. Clear filters or upload a valid Jira CSV."

    done_count = int(df["is_done"].sum()) if "is_done" in df else 0
    open_count = int(len(df) - done_count)
    blocked_df = df[df["blocked"]].sort_values(["priority", "days_open"], ascending=[True, False])
    aging_df = df[df["days_open"] > 10].sort_values("days_open", ascending=False)
    dependency_df = df[(df["has_dependency"]) | (df["blocked"])].sort_values(["blocked", "days_open"], ascending=[False, False])
    high_priority_df = df[df["priority"].str.lower().isin(["critical", "high"])].sort_values(["blocked", "days_open"], ascending=[False, False])
    defect_df = df[df["is_defect"]].sort_values(["blocked", "days_open"], ascending=[False, False])
    high_risk_df = df[df["is_high_risk"]].sort_values(["blocked", "days_open"], ascending=[False, False])

    # Greetings / help
    if q in ["help", "what can you do", "questions", "sample questions"] or "what can i ask" in q:
        return (
            "You can ask me questions grounded in the uploaded Jira data and delivery governance context. Examples:\n\n"
            "- Give me an executive summary.\n"
            "- What are the top blockers?\n"
            "- Which tickets are aging?\n"
            "- What are the dependency risks?\n"
            "- Which sprint has highest spillover?\n"
            "- Which resources are overloaded?\n"
            "- What should I discuss with the client?\n"
            "- Generate RAID actions for this sprint.\n"
            "- Explain sprint predictability / velocity / spillover."
        )

    # Executive / client summary
    if any(term in q for term in ["summary", "executive", "overall", "status report", "client update", "governance summary", "health"]):
        recs = "\n".join(f"- {r}" for r in _governance_recommendations(agg))
        return (
            f"**Executive Delivery Summary**\n\n"
            f"- Current scope: **{agg['total_tickets']} Jira item(s)** across **{agg['total_sprints']} sprint(s)**.\n"
            f"- Completion: **{done_count}/{agg['total_tickets']} tickets closed** and **{_fmt_num(agg['planned_vs_completed_pct'])}% story-point completion**.\n"
            f"- Delivery health: **{agg['health_score']}/100 ({agg['rag_status']})**.\n"
            f"- Risks: **{agg['blocked_count']} blocked**, **{agg['aging_tickets']} aging**, **{agg['dependency_count']} dependency-risk**, **{agg['defect_count']} defect(s)**.\n"
            f"- Spillover: **{_fmt_num(agg['total_spillover_story_points'], 0)} SP** remaining, **{_fmt_num(agg['spillover_pct'])}%** of planned scope.\n\n"
            f"**Recommended Governance Actions**\n{recs}"
        )

    # Blockers
    if any(term in q for term in ["blocked", "blocker", "blocking", "impediment", "stuck"]):
        return (
            f"**Blocker Analysis**\n\n"
            f"There are **{agg['blocked_count']} blocked item(s)** in the current filtered view. "
            f"Blocked items impact **{_fmt_num(blocked_df['story_points'].sum(), 0)} story point(s)**.\n\n"
            f"**Top Blocked Items**\n{_ticket_list(blocked_df, 10)}\n\n"
            f"**Client Governance Ask**\n- Confirm owner and ETA for each blocker.\n- Identify whether each blocker is internal, client-side, vendor-side or cross-team.\n- Agree escalation path for blockers impacting critical path."
        )

    # Aging
    if any(term in q for term in ["aging", "ageing", "old", "days open", "stale", "overdue"]):
        return (
            f"**Aging Ticket Analysis**\n\n"
            f"There are **{agg['aging_tickets']} item(s)** open for more than 10 days. "
            f"The oldest item is **{int(aging_df['days_open'].max()) if not aging_df.empty else 0} days** open.\n\n"
            f"**Top Aging Items**\n{_ticket_list(aging_df, 10)}\n\n"
            f"**Governance Recommendation**\n- Review aging tickets in triage.\n- Confirm whether each item should be closed, split, de-scoped, or escalated."
        )

    # Dependencies
    if any(term in q for term in ["dependency", "dependencies", "bottleneck", "external", "vendor", "cross team", "cross-team"]):
        owners = dependency_df["assignee"].nunique() if not dependency_df.empty else 0
        return (
            f"**Dependency Analysis**\n\n"
            f"The current view has **{agg['dependency_count']} dependency-risk item(s)** across **{owners} owner(s)**. "
            f"Dependency-risk items impact **{_fmt_num(dependency_df['story_points'].sum(), 0)} story point(s)**.\n\n"
            f"**Key Dependency Items**\n{_ticket_list(dependency_df, 10)}\n\n"
            f"**Client Decision Points**\n- Confirm dependency owner.\n- Confirm ETA and mitigation.\n- Escalate vendor/client-side dependencies that block sprint objectives."
        )

    # Sprint trend / velocity / predictability / spillover
    if any(term in q for term in ["sprint", "velocity", "predictability", "spillover", "planned", "completed", "commitment"]):
        if sprint_summary.empty:
            return "No sprint summary is available for the current view."
        best_velocity = sprint_summary.sort_values("velocity", ascending=False).head(1).iloc[0]
        worst_predictability = sprint_summary.sort_values("predictability", ascending=True).head(1).iloc[0]
        highest_spillover = sprint_summary.sort_values("spillover_story_points", ascending=False).head(1).iloc[0]
        lines = []
        for _, r in sprint_summary.head(10).iterrows():
            lines.append(
                f"- **{r['sprint']}**: Planned {_fmt_num(r['planned_story_points'],0)} SP, "
                f"Completed {_fmt_num(r['completed_story_points'],0)} SP, "
                f"Predictability {_fmt_num(r['predictability'])}%, Spillover {_fmt_num(r['spillover_story_points'],0)} SP"
            )
        return (
            f"**Sprint Performance Analysis**\n\n"
            f"- Average velocity: **{_fmt_num(agg['average_velocity'])} SP**.\n"
            f"- Average predictability: **{_fmt_num(agg['average_predictability'])}%**.\n"
            f"- Highest velocity sprint: **{best_velocity['sprint']}** with **{_fmt_num(best_velocity['velocity'],0)} SP**.\n"
            f"- Lowest predictability sprint: **{worst_predictability['sprint']}** at **{_fmt_num(worst_predictability['predictability'])}%**.\n"
            f"- Highest spillover sprint: **{highest_spillover['sprint']}** with **{_fmt_num(highest_spillover['spillover_story_points'],0)} SP**.\n\n"
            f"**Sprint Facts**\n" + "\n".join(lines)
        )

    # Resource/capacity
    if any(term in q for term in ["resource", "assignee", "owner", "capacity", "overload", "overloaded", "workload", "utilization", "utilisation"]):
        if resource_workload.empty:
            return "No resource workload data is available in the current view."
        avg = resource_workload["workload_story_points"].mean()
        overloaded = resource_workload[resource_workload["workload_story_points"] >= 1.5 * avg]
        top_rows = []
        for _, r in resource_workload.head(10).iterrows():
            top_rows.append(
                f"- **{r['assignee']}**: {_fmt_num(r['workload_story_points'],0)} SP, "
                f"{int(r['ticket_count'])} ticket(s), {int(r['open_tickets'])} open, {int(r['blocked_tickets'])} blocked"
            )
        overload_text = "No assignee is above 150% of average workload." if overloaded.empty else ", ".join(overloaded["assignee"].astype(str).tolist())
        return (
            f"**Resource Workload Analysis**\n\n"
            f"- Active resources: **{agg['resources_count']}**.\n"
            f"- Average workload: **{_fmt_num(avg)} SP per resource**.\n"
            f"- Overload check: **{overload_text}**\n\n"
            f"**Workload by Assignee**\n" + "\n".join(top_rows) +
            "\n\n**Governance Recommendation**\n- Rebalance blockers and high-priority work from overloaded owners.\n- Confirm capacity before accepting additional scope."
        )

    # Quality / defects
    if any(term in q for term in ["defect", "bug", "quality", "qa", "testing", "test"]):
        return (
            f"**Quality & Defect Analysis**\n\n"
            f"- Defects: **{agg['defect_count']}**.\n"
            f"- Defect rate: **{_fmt_num(agg['defect_rate'])}%** of total Jira items.\n"
            f"- Quality status: **{'Needs Attention' if agg['defect_rate'] > 10 else 'Within Target'}**.\n\n"
            f"**Top Defect Items**\n{_ticket_list(defect_df, 10)}\n\n"
            f"**Quality Governance Ask**\n- Prioritize critical/high defects.\n- Confirm defect owner and target fix sprint.\n- Run root-cause review if defect rate stays above 10%."
        )

    # RAID / actions / next steps
    if any(term in q for term in ["raid", "action", "actions", "next step", "next steps", "decision", "decisions", "mitigation", "recommend", "recommendation", "client pack", "client meeting"]):
        raid_items = []
        for _, r in blocked_df.head(5).iterrows():
            severity = "Critical" if str(r.get("priority", "")).lower() == "critical" or int(r.get("days_open", 0)) > 10 else "High"
            raid_items.append(f"- **Issue | {r['ticket_key']} | {severity}**: Blocked {r.get('priority')} item. Action: Confirm owner, ETA and escalation path.")
        for _, r in dependency_df.head(3).iterrows():
            dep = str(r.get("dependency_text", "")).strip() or r.get("ticket_key", "")
            raid_items.append(f"- **Dependency | {dep} | High**: Linked dependency in current sprint. Action: Agree dependency owner and due date.")
        if not raid_items:
            raid_items.append("- No major RAID items found in the current filter context.")
        recs = "\n".join(f"- {r}" for r in _governance_recommendations(agg))
        return (
            f"**RAID / Action Tracker Draft**\n\n" + "\n".join(raid_items[:8]) +
            f"\n\n**Recommended Next Steps**\n{recs}"
        )

    # Priority / risks
    if any(term in q for term in ["risk", "risks", "critical", "high priority", "priority", "red", "amber", "rag"]):
        return (
            f"**Risk Analysis**\n\n"
            f"- RAG status: **{agg['rag_status']}** with health score **{agg['health_score']}/100**.\n"
            f"- High-risk items: **{agg['high_risk_count']}**.\n"
            f"- Critical/High priority items: **{len(high_priority_df)}**.\n"
            f"- Blocked: **{agg['blocked_count']}**, Aging: **{agg['aging_tickets']}**, Dependencies: **{agg['dependency_count']}**.\n\n"
            f"**Top Risk Items**\n{_ticket_list(high_risk_df, 10)}\n\n"
            f"**Governance Recommendation**\n- Discuss top risks in the client meeting and capture owner/date/action in the RAID tracker."
        )

    # Data exploration fallback by keyword: try match ticket/sprint/status/assignee/priority/type
    tokens = [t.strip().lower() for t in q.replace(",", " ").split() if len(t.strip()) >= 3]
    mask = pd.Series(False, index=df.index)
    searchable_cols = ["ticket_key", "sprint", "status", "priority", "assignee", "issue_type", "dependency_text", "risk_category"]
    for col in searchable_cols:
        if col in df.columns:
            text_col = df[col].fillna("").astype(str).str.lower()
            for token in tokens:
                mask = mask | text_col.str.contains(token, na=False, regex=False)
    matched = df[mask]
    if not matched.empty:
        return (
            f"I found **{len(matched)} matching Jira item(s)** for your question in the current filtered view.\n\n"
            f"{_ticket_list(matched.sort_values(['blocked', 'days_open'], ascending=[False, False]), 10)}"
        )

    return (
        "I could not map the question to a specific Jira/governance insight. Try asking about: "
        "executive summary, blockers, aging tickets, dependency risks, sprint velocity, predictability, spillover, resource workload, defects, RAID actions, or client next steps."
    )
# ============================================================
# MAIN APP
# ============================================================
def main():
    apply_app_style()
    st.markdown("<span id='app-top' class='app-top-anchor'></span>", unsafe_allow_html=True)
    st.markdown("<a class='smooth-top-button' href='#app-top' title='Smoothly Scroll To Top'>↑</a>", unsafe_allow_html=True)
    init_filter_state()

    now = datetime.now(ZoneInfo("Asia/Kolkata"))
    greeting = "Good Morning" if now.hour < 12 else "Good Afternoon" if now.hour < 18 else "Good Evening"
    current_datetime = now.strftime("%A, %B %d %Y · %I:%M %p IST")

    st.markdown(
        f"""
        <div class='hero'>
          <div class='hero-top'>
            <div>
              <div class='hero-badge'>AI Agent - Delivery Manager</div>
              <h1>{greeting}, Delivery Manager</h1>
              <p>Executive Delivery Governance Dashboard For Sprint Health, Risks, Dependencies, Resource Workload And KPI Reporting.</p>
            </div>
            <div class='hero-meta'>
              <strong>{current_datetime}</strong><br/>
              Executive Delivery View Updated In Real Time.
            </div>
          </div>
        </div>
        """,
        unsafe_allow_html=True,
    )

    st.markdown("<div class='upload-panel'>", unsafe_allow_html=True)
    st.markdown("### Upload Jira CSV or Download Template")
    st.caption("Use your Jira export CSV. If the file is not ready, download the template below and populate it with sprint delivery data.")

    template_col, upload_col = st.columns([1, 2], gap="large")
    with template_col:
        st.download_button(
            label="Download Jira CSV Template",
            data=build_jira_template_csv_bytes(),
            file_name="jira_delivery_governance_template.csv",
            mime="text/csv",
            use_container_width=True,
            help="Downloads a Jira CSV template with sample rows and all required columns for the dashboard.",
        )
        st.download_button(
            label="Download Blank Template",
            data=build_blank_jira_template_csv_bytes(),
            file_name="jira_delivery_governance_blank_template.csv",
            mime="text/csv",
            use_container_width=True,
            help="Downloads only the column headers if you want to start from a blank file.",
        )
        with st.expander("Template column guide"):
            st.markdown(
                """
                **Recommended columns:** Key, Sprint, Type, Status, Assignee, Priority, Story Points, Days Open, Dependency, Created, Updated, Resolved, Blocked, Labels.

                **Accepted values:** Status can be Done, In Progress, Open, Blocked. Type can be Story, Bug, Task. Blocked can be Yes/No.
                """
            )
    with upload_col:
        uploaded_file = st.file_uploader(
            "Upload Jira CSV",
            type=["csv"],
            help="Upload Jira export CSV. The dashboard and PPT will use the uploaded data.",
        )
        if uploaded_file is None:
            st.info("Upload a Jira CSV file to view the Delivery Manager AI Agent dashboard.")
    st.markdown("</div>", unsafe_allow_html=True)

    if uploaded_file is None:
        return

    raw_df = load_data(uploaded_file)
    if raw_df.empty:
        st.error("The uploaded CSV is empty or could not be read.")
        return

    base_df = prepare_model(raw_df)

    # Sidebar filters
    with st.sidebar:
        st.header("Filters")
        if st.button("Clear All Filters", use_container_width=True):
            clear_filters()
            st.rerun()
        st.selectbox("Sprint", get_options(base_df["sprint"]), key="f_sprint")
        st.selectbox("Status", get_options(base_df["status"]), key="f_status")
        st.selectbox("Priority", get_options(base_df["priority"]), key="f_priority")
        st.selectbox("Assignee", get_options(base_df["assignee"]), key="f_assignee")
        st.selectbox("Risk Category", get_options(base_df["risk_category"]), key="f_risk_category")
        st.selectbox("Month", get_options(base_df["month"]), key="f_month")
        st.selectbox("Week", get_options(base_df["week"]), key="f_week")
        st.caption("Filters apply to all visuals, AI answers and PPT export.")

    filtered_df = apply_filters(base_df)
    metrics = compute_agile_metrics(filtered_df)
    model_df = metrics["df"]
    agg = metrics["agg_metrics"]
    sprint_summary = metrics["sprint_summary"]
    resource_workload = metrics["resource_workload"]
    defect_trend = metrics["defect_trend"]

    st.markdown(f"<div class='filter-note'>Current Filter Context: {selected_filters_text()}</div>", unsafe_allow_html=True)

    if model_df.empty:
        st.warning("No records found for the selected filter combination. Click Clear All Filters to reset.")
        return

    st.markdown("<div class='section-title'>📊 Executive KPI Snapshot</div>", unsafe_allow_html=True)
    render_kpi_grid(agg)

    # Delivery health charts
    st.markdown("<div class='section-title'>🚦 Delivery Health</div>", unsafe_allow_html=True)
    col1, col2 = st.columns(2, gap="large")
    with col1:
        render_chart(make_status_chart(model_df), "chart_status")
        st.markdown("<div class='hint'>Use the Status filter in the sidebar to update all visuals.</div>", unsafe_allow_html=True)
    with col2:
        render_chart(make_risk_chart(model_df), "chart_risk")
        st.markdown("<div class='hint'>Use the Risk Category filter in the sidebar to update all visuals.</div>", unsafe_allow_html=True)

    # Dependency and resource charts
    st.markdown("<div class='section-title'>🔗 Dependency & Resource Analysis</div>", unsafe_allow_html=True)
    col3, col4 = st.columns(2, gap="large")
    with col3:
        render_chart(make_dependency_chart(model_df), "chart_dependency")
        st.markdown("<div class='hint'>Use the Assignee filter in the sidebar to review dependency ownership.</div>", unsafe_allow_html=True)
    with col4:
        render_chart(make_resource_chart(resource_workload), "chart_resource")
        st.markdown("<div class='hint'>Use the Assignee filter in the sidebar to focus resource workload and sprint facts.</div>", unsafe_allow_html=True)

    # Agile metrics charts
    st.markdown("<div class='section-title'>📈 Agile Metrics</div>", unsafe_allow_html=True)
    col5, col6 = st.columns(2, gap="large")
    with col5:
        render_chart(make_velocity_chart(sprint_summary), "chart_velocity")
        st.markdown("<div class='hint'>Use the Sprint filter in the sidebar to focus all visuals on one sprint.</div>", unsafe_allow_html=True)
    with col6:
        render_chart(make_defect_trend(defect_trend, sprint_summary), "chart_defect")
        st.markdown("<div class='hint'>Use sidebar filters to inspect defect trend by sprint, priority or assignee.</div>", unsafe_allow_html=True)

    col7, col8 = st.columns(2, gap="large")
    with col7:
        render_chart(make_planned_completed_chart(sprint_summary), "chart_planned_completed")
        st.markdown("<div class='hint'>Use the Sprint filter to compare planned versus completed facts for one sprint.</div>", unsafe_allow_html=True)
    with col8:
        render_quality_kpi(agg)

    # Priority chart and data table
    st.markdown("<div class='section-title'>⚠ Priority & Risk Items</div>", unsafe_allow_html=True)
    col9, col10 = st.columns([1, 1], gap="large")
    with col9:
        render_chart(make_priority_chart(model_df), "chart_priority")
        st.markdown("<div class='hint'>Use the Priority filter in the sidebar to update all related visuals.</div>", unsafe_allow_html=True)
    with col10:
        risk_display_df = (
            model_df[["ticket_key", "sprint", "status", "priority", "assignee", "story_points", "days_open", "risk_category"]]
            .sort_values(["risk_category", "days_open"], ascending=[True, False])
            .head(20)
        )
        st.dataframe(
            display_table(risk_display_df),
            use_container_width=True,
            height=CHART_HEIGHT,
        )

    # AI Assistant and PPT
    st.markdown("<div class='section-title'>🤖 AI Assistant & Client Pack</div>", unsafe_allow_html=True)
    st.caption(
        "Ask questions about the uploaded Jira data, sprint health, blockers, dependencies, resource workload, quality, RAID actions, or governance recommendations."
    )

    suggested_questions = [
        "Give me an executive summary",
        "What are the top blockers?",
        "Which tickets are aging?",
        "What are the dependency risks?",
        "Which sprint has highest spillover?",
        "Which resources are overloaded?",
        "Generate RAID actions for the client meeting",
        "What should I discuss with the client?",
    ]

    if "ai_question" not in st.session_state:
        st.session_state["ai_question"] = ""

    sq1, sq2, sq3, sq4 = st.columns(4)
    for col, label in zip([sq1, sq2, sq3, sq4], suggested_questions[:4]):
        with col:
            if st.button(label, use_container_width=True):
                st.session_state["ai_question"] = label
    sq5, sq6, sq7, sq8 = st.columns(4)
    for col, label in zip([sq5, sq6, sq7, sq8], suggested_questions[4:]):
        with col:
            if st.button(label, use_container_width=True):
                st.session_state["ai_question"] = label

    q = st.text_input(
        "Ask the AI Assistant",
        key="ai_question",
        placeholder="Example: Prepare client talking points for blockers, dependencies and sprint recovery",
    )
    if q:
        st.markdown(answer_ai_question(q, metrics))

    st.markdown("---")
    st.subheader("Client Pack Export")
    st.caption("The PPT uses the same filtered Jira view currently shown in the dashboard.")
    pptx_bytes = generate_client_engagement_ppt(model_df)
    st.download_button(
        "📊 Download Client Engagement KPI PPT",
        data=pptx_bytes,
        file_name="client_engagement_delivery_governance_kpi_pack.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        use_container_width=True,
    )

    with st.expander("Raw Filtered Jira Data", expanded=False):
        st.dataframe(display_table(model_df), use_container_width=True)


if __name__ == "__main__":
    main()
